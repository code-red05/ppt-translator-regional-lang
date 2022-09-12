import sys
from pptx import Presentation
from translator import eng_to_hin, eng_to_tam, eng_to_mar

prs = Presentation(sys.argv[1])

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if sys.argv[3] == "hin":
                    run.text = eng_to_hin(run.text)
                elif sys.argv[3] == "tam":
                    run.text = eng_to_tam(run.text)
                elif sys.argv[3] == "mar":
                    run.text = eng_to_mar(run.text)
                else :
                    print("only 3 lang supported: hin, tam, mar")

prs.save(sys.argv[2])
